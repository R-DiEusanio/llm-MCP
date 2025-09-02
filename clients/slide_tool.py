# clients/slide_tool.py
from __future__ import annotations
import io, json, re
from typing import List
from pydantic import BaseModel, Field
from langchain_openai import ChatOpenAI
from pptx import Presentation
from pptx.util import Pt

class Slide(BaseModel):
    title: str
    bullets: List[str] = Field(default_factory=list)

class SlideDeck(BaseModel):
    subject: str
    topic: str
    slides: List[Slide]

def _draft_slides(subject: str, topic: str, n_slides: int) -> SlideDeck:
    """Chiede all'LLM un outline JSON con titoli + bullet."""
    llm = ChatOpenAI(model="gpt-4o", temperature=0.3)
    prompt = f"""
Sei un docente delle scuole superiori italiane.
Crea un piano di {n_slides} slide per la materia "{subject}" con argomento "{topic}".
Rispondi SOLO con JSON nel formato:
{{"slides":[{{"title":"...","bullets":["...","..."]}}, ...]}}
- Max 5 bullet per slide
- Linguaggio semplice e didattico
- Niente markdown
"""
    text = llm.invoke(prompt).content
    m = re.search(r"\{[\s\S]*\}", text)
    payload = json.loads(m.group(0) if m else text)
    slides = [Slide(**s) for s in payload["slides"]][:max(1, n_slides)]
    return SlideDeck(subject=subject, topic=topic, slides=slides)

def _build_pptx(deck: SlideDeck) -> io.BytesIO:
    prs = Presentation()

    # Copertina
    title_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = deck.topic
    cover.placeholders[1].text = deck.subject

    # Slide contenuto
    content_layout = prs.slide_layouts[1]
    for s in deck.slides:
        sl = prs.slides.add_slide(content_layout)
        sl.shapes.title.text = s.title
        tf = sl.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(s.bullets):
            if i == 0:
                tf.text = b
                tf.paragraphs[0].font.size = Pt(20)
            else:
                p = tf.add_paragraph()
                p.text = b
                p.level = 1
                p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def generate_slides_pptx(subject: str, topic: str, n_slides: int = 10):
    """Ritorna un BytesIO del PPTX generato."""
    deck = _draft_slides(subject, topic, n_slides)
    return _build_pptx(deck)
